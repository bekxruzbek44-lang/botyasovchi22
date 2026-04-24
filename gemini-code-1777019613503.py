import os
import asyncio
import json
import logging
from aiogram import Bot, Dispatcher, types, F
from aiogram.filters import Command
from pptx import Presentation
import google.generativeai as genai

# Tokenlarni xavfsiz yuklash (Render sozlamalaridan oladi)
BOT_TOKEN = os.getenv("BOT_TOKEN")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

# Bot va Gemini sozlamalari
bot = Bot(token=BOT_TOKEN)
dp = Dispatcher()
genai.configure(api_key=GEMINI_API_KEY)

logging.basicConfig(level=logging.INFO)

def create_pptx(topic):
    model = genai.GenerativeModel('gemini-pro')
    # Prezentatsiya uchun prompt
    prompt = f"'{topic}' mavzusida 5 ta slaydli prezentatsiya tuz. Har bir slayd uchun 'title' va 'body' (matn) qismlarini JSON formatida qaytar."
    response = model.generate_content(prompt)
    
    # JSONni tozalash va o'qish
    text = response.text.replace('```json', '').replace('```', '')
    data = json.loads(text)
    
    # PPTX faylni yaratish
    prs = Presentation()
    for item in data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = item.get('title', 'Sarlavha')
        slide.placeholders[1].text = item.get('body', 'Matn')
        
    file_name = f"{topic.replace(' ', '_')}.pptx"
    prs.save(file_name)
    return file_name

@dp.message(Command("start"))
async def start(message: types.Message):
    await message.answer("Salom! Mavzuni yozing, men sizga avtomatik prezentatsiya tayyorlab beraman.")

@dp.message()
async def generate(message: types.Message):
    await message.answer("⏳ Prezentatsiya tayyorlanmoqda, iltimos kuting...")
    try:
        file_path = create_pptx(message.text)
        await message.answer_document(types.FSInputFile(file_path))
        os.remove(file_path) # Faylni yuborgach o'chirish
    except Exception as e:
        await message.answer(f"Xatolik yuz berdi: {e}")

async def main():
    await dp.start_polling(bot)

if __name__ == "__main__":
    asyncio.run(main())